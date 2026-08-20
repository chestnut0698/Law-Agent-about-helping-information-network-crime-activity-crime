import json
from pathlib import Path
import pymupdf
import fitz  # PyMuPDF，PDF 文本提取
from app.config import WORKSPACE_DIR

# ---------- 全局变量 ----------
_ocr_engine = None

# ---------- 文件类型分类 ----------
TEXT_EXTS = {
    ".txt", ".md", ".csv", ".json", ".xml", ".py", ".js",
    ".html", ".css", ".yaml", ".yml", ".ini", ".cfg", ".log",
    ".rtf", ".tex",
}
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff", ".tif"}
PDF_EXTS = {".pdf"}
DOCX_EXTS = {".docx"}          # 新版 Word（Office Open XML）
DOC_EXTS = {".doc"}            # 旧版 Word（需要 antiword 或 libreoffice）
PPTX_EXTS = {".pptx"}          # 新版 PowerPoint
XLSX_EXTS = {".xlsx", ".xls"}  # Excel

# 汇总，方便判断
ALL_SUPPORTED = TEXT_EXTS | IMAGE_EXTS | PDF_EXTS | DOCX_EXTS | DOC_EXTS | PPTX_EXTS | XLSX_EXTS


# ---------- 路径安全工具 ----------

def _safe_resolve(task_id: str, sub_path: str, filename: str) -> Path | None:
    """拼接并校验路径，防止穿越攻击。返回 None 表示非法。"""
    rel = Path(sub_path.strip("/").replace("\\", "/")) / filename if sub_path else Path(filename)
    target = (WORKSPACE_DIR / task_id / rel).resolve()
    if not str(target).startswith(str(WORKSPACE_DIR.resolve())):
        return None
    return target


# ---------- OCR 引擎（懒加载） ----------

def _get_ocr():
    global _ocr_engine
    if _ocr_engine is None:
        from paddleocr import PaddleOCR
        _ocr_engine = PaddleOCR(
            use_angle_cls=True,
            lang="ch",
            enable_mkldnn=False,   # ← 关键：绕过 PIR/OneDNN 崩溃
        )
    return _ocr_engine


def _ocr_image_bytes(image_bytes: bytes) -> str:
    import numpy as np
    from PIL import Image
    import io
    engine = _get_ocr()
    img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
    try:
        result = engine.predict(np.array(img))  # 3.x
        lines = []
        for block in result:
            if isinstance(block, dict):
                lines.extend(block.get("rec_texts", block.get("texts", [])))
            elif isinstance(block, list):
                lines.extend(line[1][0] for line in block)
    except AttributeError:
        result = engine.ocr(np.array(img), cls=True)  # 2.x
        lines = [line[1][0] for block in result for line in (block or [])]
    return "\n".join(lines)



# ---------- 文本解码（多编码回退） ----------

def _read_text_file(path: Path) -> str:
    """尝试多种编码读取文本文件。"""
    for enc in ("utf-8", "gbk", "gb2312", "latin-1"):
        try:
            return path.read_text(encoding=enc)
        except UnicodeDecodeError:
            continue
    # 最后兜底：忽略错误
    return path.read_text(encoding="utf-8", errors="ignore")


# ---------- 各格式提取函数 ----------

def _extract_pdf(path: Path) -> str:
    """PDF：先抽文本，空页再 OCR。"""
    text_parts = []
    doc = fitz.open(path)
    try:
        for page in doc:
            page_text = page.get_text("text") or ""
            if page_text.strip():
                text_parts.append(page_text)
            else:
                # 扫描页 → 渲染 → OCR
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                text_parts.append(_ocr_image_bytes(pix.tobytes("png")))
    finally:
        doc.close()
    return "\n\n".join(text_parts)


def _extract_docx(path: Path) -> str:
    """Word 2007+ (.docx)：python-docx 提取段落 + 表格。"""
    from docx import Document
    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_doc(path: Path) -> str:
    """
    旧版 Word (.doc)：优先用 libreoffice 转 docx 再提取；
    其次尝试 antiword；都失败则返回提示。
    """
    import subprocess, tempfile, os

    # 方法1：libreoffice 转 docx
    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "docx",
                 "--outdir", tmpdir, str(path)],
                check=True, capture_output=True, timeout=30,
            )
            converted = Path(tmpdir) / (path.stem + ".docx")
            if converted.exists():
                return _extract_docx(converted)
        except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
            pass

    # 方法2：antiword
    try:
        result = subprocess.run(
            ["antiword", str(path)], check=True, capture_output=True, timeout=15,
        )
        return result.stdout.decode("utf-8", errors="ignore")
    except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
        pass

    return "[无法提取 .doc 文件内容：请安装 libreoffice 或 antiword]"


def _extract_pptx(path: Path) -> str:
    """PowerPoint (.pptx)：提取每页文本框内容。"""
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for idx, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            parts.append(f"[Slide {idx}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    """Excel (.xlsx/.xls)：逐 sheet 提取，输出为 CSV 风格文本。"""
    from openpyxl import load_workbook
    wb = load_workbook(str(path), data_only=True, read_only=True)
    parts = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            rows.append("\t".join(cells))
        if rows:
            parts.append(f"[Sheet: {ws.title}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


# ---------- 对外工具函数 ----------

def list_user_files(task_id: str) -> str:
    """
    递归列出会话目录下的所有文件（含子目录）。
    返回 JSON 字符串，每项含 name / relative_path / size / suffix。
    """
    target = WORKSPACE_DIR / task_id
    if not target.is_dir():
        return json.dumps({"error": f"未找到目录: {target}"}, ensure_ascii=False)

    try:
        files = []
        for entry in sorted(target.rglob("*")):
            if entry.is_file():
                rel = entry.relative_to(target)
                files.append({
                    "name": entry.name,
                    "relative_path": str(rel).replace("\\", "/"),
                    "size": entry.stat().st_size,
                    "suffix": entry.suffix.lower(),
                    "supported": entry.suffix.lower() in ALL_SUPPORTED,
                })
        return json.dumps({"files": files, "count": len(files)},
                         ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def read_any_file(task_id: str, filename: str, sub_path: str = "") -> str:
    """
    统一读取入口：文本/图片/PDF/Word 都返回纯文本。
    - sub_path: 可选，子目录相对路径，如 "证据/图片/"
    """
    # ★ 关键修复：定义 rel 变量
    rel = Path(sub_path) / filename if sub_path else Path(filename)
    target = (WORKSPACE_DIR / task_id / rel).resolve()

    # 路径穿越防护
    if not str(target).startswith(str(WORKSPACE_DIR.resolve())):
        return json.dumps({"error": "Access denied"}, ensure_ascii=False)
    if not target.is_file():
        return json.dumps({"error": f"File not found: {rel}"}, ensure_ascii=False)

    ext = target.suffix.lower()
    try:
        # ---- 文本文件 ----
        if ext in TEXT_EXTS:
            for enc in ("utf-8", "gbk", "gb2312", "latin-1"):
                try:
                    content = target.read_text(encoding=enc)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                return json.dumps({"error": "Cannot decode text file"})
            return json.dumps({
                "filename": str(rel).replace("\\", "/"),
                "type": "text",
                "text": content
            }, ensure_ascii=False, indent=2)

        # ---- PDF ----
        elif ext in PDF_EXTS:
            try:
                doc = fitz.open(target)
            except Exception as e:
                return json.dumps({"error": f"PDF 打开失败: {e}"}, ensure_ascii=False)
            text_parts = []
            try:
                for page in doc:
                    page_text = page.get_text("text") or ""
                    if page_text.strip():
                        text_parts.append(page_text)
                    else:
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        text_parts.append(_ocr_image_bytes(pix.tobytes("png")))
            except Exception as e:
                doc.close()
                return json.dumps({"error": f"PDF 读取失败: {e}"}, ensure_ascii=False)
            doc.close()
            return json.dumps({
                "filename": str(rel).replace("\\", "/"),
                "type": "pdf",
                "text": "\n\n".join(text_parts)
            }, ensure_ascii=False, indent=2)

        # ---- DOCX ----
        elif ext in DOCX_EXTS:
            try:
                from docx import Document
                doc = Document(str(target))
                parts = [p.text for p in doc.paragraphs if p.text]
                for table in doc.tables:
                    rows = ["\t".join(cell.text.strip() for cell in row.cells if cell.text.strip()) for row in table.rows]
                    parts.append("\n".join(rows))
                text = "\n".join(parts)
            except ImportError:
                return json.dumps({"error": "python-docx 未安装，无法读取 .docx 文件"}, ensure_ascii=False)
            except Exception as e:
                return json.dumps({"error": f"DOCX 读取失败: {e}"}, ensure_ascii=False)
            return json.dumps({
                "filename": str(rel).replace("\\", "/"),
                "type": "docx",
                "text": text
            }, ensure_ascii=False, indent=2)

        # ---- DOC（旧版 Word，需 libreoffice 转换）----
        elif ext in DOC_EXTS:
            try:
                import subprocess
                import tempfile
                # 转换为 docx
                with tempfile.TemporaryDirectory() as tmpdir:
                    subprocess.run(
                        ["libreoffice", "--headless", "--convert-to", "docx", str(target), "--outdir", tmpdir],
                        capture_output=True, timeout=60
                    )
                    converted = Path(tmpdir) / (target.stem + ".docx")
                    if not converted.exists():
                        return json.dumps({"error": "DOC 转换失败"}, ensure_ascii=False)
                    from docx import Document
                    doc = Document(str(converted))
                    parts = [p.text for p in doc.paragraphs if p.text]
                    text = "\n".join(parts)
            except FileNotFoundError:
                return json.dumps({"error": "未安装 libreoffice，无法读取 .doc 文件"}, ensure_ascii=False)
            except Exception as e:
                return json.dumps({"error": f"DOC 读取失败: {e}"}, ensure_ascii=False)
            return json.dumps({
                "filename": str(rel).replace("\\", "/"),
                "type": "doc",
                "text": text
            }, ensure_ascii=False, indent=2)

        # ---- 图片 OCR ----
        elif ext in IMAGE_EXTS:
            text = _ocr_image_bytes(target.read_bytes())
            return json.dumps({
                "filename": str(rel).replace("\\", "/"),
                "type": "image_ocr",
                "text": text
            }, ensure_ascii=False, indent=2)

        # ---- PPTX ----
        elif ext in PPTX_EXTS:
            try:
                from pptx import Presentation
                prs = Presentation(str(target))
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text)
                text = "\n".join(parts)
            except ImportError:
                return json.dumps({"error": "python-pptx 未安装，无法读取 .pptx 文件"}, ensure_ascii=False)
            except Exception as e:
                return json.dumps({"error": f"PPTX 读取失败: {e}"}, ensure_ascii=False)
            return json.dumps({
                "filename": str(rel).replace("\\", "/"),
                "type": "pptx",
                "text": text
            }, ensure_ascii=False, indent=2)

        # ---- XLSX/XLS ----
        elif ext in XLSX_EXTS:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(target), data_only=True)
                parts = []
                for sheet in wb.worksheets:
                    sheet_lines = []
                    for row in sheet.iter_rows(values_only=True):
                        cells = [str(c) if c is not None else "" for c in row]
                        sheet_lines.append("\t".join(cells))
                    parts.append(f"=== Sheet: {sheet.title} ===\n" + "\n".join(sheet_lines))
                text = "\n\n".join(parts)
            except ImportError:
                return json.dumps({"error": "openpyxl 未安装，无法读取 .xlsx/.xls 文件"}, ensure_ascii=False)
            except Exception as e:
                return json.dumps({"error": f"Excel 读取失败: {e}"}, ensure_ascii=False)
            return json.dumps({
                "filename": str(rel).replace("\\", "/"),
                "type": "excel",
                "text": text
            }, ensure_ascii=False, indent=2)

        else:
            return json.dumps({"error": f"不支持的文件类型: {ext}"})

    except Exception as e:
        return json.dumps({"error": f"读取失败: {e}"})
